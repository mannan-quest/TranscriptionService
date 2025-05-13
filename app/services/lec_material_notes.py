import asyncio
import os
from typing import List
import PyPDF2
from openai import OpenAI
import pptx
from spire.doc import *
from spire.doc.common import *

from fastapi import UploadFile
from supabase import create_client

from app.services.translation_service import TranslationAnalysisService
from app.core.config import settings

class LectureMaterialNotes:
    def __init__(self, lecture_material_id, file_path, filetype):
        self.lecture_material_id = lecture_material_id
        self.file_path = file_path
        self.filetype = filetype
        self.SUPABASE_URL = os.getenv("SUPABASE_URL")
        self.SUPABASE_KEY = os.getenv("SUPABASE_KEY")
        self.client = OpenAI(api_key=settings.OPENAI_API_KEY)
        self.supabase = create_client(self.SUPABASE_URL, self.SUPABASE_KEY)
        self.translation_service = TranslationAnalysisService()

    def analyze_material(self):
        # Analyze the lecture material and route the material to the appropriate function
        # based on the type of material (e.g., txt, pdf, docx, etc.)
        self._uploadtoVectorStore()  # Upload to OpenAI Vector Store
        if self.filetype == 'pdf':
            return self.process_pdf()
        elif self.filetype == 'ppt' or self.filetype == 'pptx':
            return self.process_pptx()
        elif self.filetype == 'doc' or self.filetype == 'docx':
            return self.process_docx()
        elif self.filetype == 'txt':
            return self.process_txt()
        else:
            self.update_progress(1.0) # 100% done
            raise Exception(f"Unsupported file type: {self.filetype}")

    def update_progress(self, value: float):
        self.supabase.table("lecture_materials") \
            .update({"progress": value}) \
            .eq("material_id", self.lecture_material_id) \
            .execute()

    def _uploadtoVectorStore(self):
        """Upload the file to OpenAI Vector Store for this specific lecture"""
        # 1) Get the vectorstore id from supabase
        # 2) Update the vectorstore with the file ID
        lecture_id = self.supabase.table("lecture_materials").select("lecture_id").eq("material_id", self.lecture_material_id).execute()
        lecture_id = lecture_id.data[0]['lecture_id']
        vectorstore_id = self.supabase.table("lectures").select("vectorstore_id").eq("lecture_id", lecture_id).execute()
        vectorstore_id = vectorstore_id.data[0]['vectorstore_id']

        with open(self.file_path, 'rb') as file:
            file_response = self.client.vector_stores.files.upload_and_poll(vector_store_id=vectorstore_id, file=file)
            file_id = file_response.id
            print(f"File uploaded to OpenAI Vector Store with ID: {file_id}")

    async def process_txt(self):
        try:
            # Extract text from a TXT file
            with open(self.file_path, 'r') as file:
                text = file.read()
                self.update_progress(0.2)  # 20% done
                text = sanitize_text(text)
                self.update_progress(0.5)  # 50% done
                analysis = await self.translation_service.analyze_material_text([text])
                self.update_progress(0.7)
                self.supabase.table("lecture_materials").update({
                    "notes": analysis,
                    "paragraphs": [text]
                }).eq("material_id", self.lecture_material_id).execute()
                self.update_progress(1.0)  # 100% done
                print(f"TXT processing completed for material {self.lecture_material_id}")
        except Exception as e:
            # If something fails, update the DB accordingly
            error_msg = str(e)
            try:
                # Update the database with the error status
                self.supabase.table("lecture_materials").update({
                    "error": sanitize_text(error_msg),
                    "progress": 0.0  # Reset progress
                }).eq("material_id", self.lecture_material_id).execute()
            except Exception as db_error:
                print(f"Error updating database with error status: {db_error}")
                
            print(f"Error processing TXT material {self.lecture_material_id}: {error_msg}")
            raise e
        finally:
            # Clean up temporary file
            try:
                if os.path.exists(self.file_path):
                    os.remove(self.file_path)
                    print(f"Temporary file removed: {self.file_path}")
            except Exception as cleanup_error:
                print(f"Error cleaning up temporary file: {cleanup_error}")
                        
    async def process_docx(self):
        try:
            # Extract text from a DOCX or DOC file
            progress = 0.2
            self.update_progress(progress)  # 20% done
            doc = Document()
            doc.LoadFromFile(self.file_path)
            doc_text = sanitize_text(doc.GetText())
            doc.Close()
            
            self.update_progress(0.5)  # 50% done
    
            analysis = await self.translation_service.analyze_material_text([doc_text])
    
            self.update_progress(0.7)  # 70% done

            # Insert notes
            print("Inserting notes")
            self.supabase.table("lecture_materials").update({
                "notes": analysis, 
                "paragraphs": [doc_text]
            }).eq("material_id", self.lecture_material_id).execute()
    
            # Mark done
            self.update_progress(1.0)  # 100% done
            print(f"DOCX processing completed for material {self.lecture_material_id}")
     
        except Exception as e:
            # If something fails, update the DB accordingly
            error_msg = str(e)
            try:
                # Update the database with the error status
                self.supabase.table("lecture_materials").update({
                    "error": sanitize_text(error_msg),
                    "progress": 0.0  # Reset progress
                }).eq("material_id", self.lecture_material_id).execute()
            except Exception as db_error:
                print(f"Error updating database with error status: {db_error}")
                
            print(f"Error processing DOCX material {self.lecture_material_id}: {error_msg}")
            raise e  # Re-raise the exception to propagate it
        finally:
            # Clean up temporary file
            try:
                if os.path.exists(self.file_path):
                    os.remove(self.file_path)
                    print(f"Temporary file removed: {self.file_path}")
            except Exception as cleanup_error:
                print(f"Error cleaning up temporary file: {cleanup_error}")   

     
    async def process_pptx(self):
        """Process PowerPoint files with proper error handling."""
        try:
            # Extract text from a PPTX file
            prs = pptx.Presentation(self.file_path)
            text_runs = []
            progress = 0.2
    
            self.update_progress(progress)  # 20% done
    
            total_slides = len(prs.slides)
            if total_slides == 0:
                raise Exception("No slides found in PPTX")
            
            progress_step = 0.5 / total_slides
    
            # Read text from each slide
            for slide in prs.slides:
                try:
                    progress += progress_step
                    self.update_progress(progress)
                    slide_text = ""
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_text += run.text
                    
                    # Only add non-empty slide text
                    if slide_text.strip():
                        text_runs.append(sanitize_text(slide_text))
                except Exception as slide_error:
                    print(f"Error processing slide: {slide_error}")
                    # Continue with next slide instead of failing completely
                    continue
            
            if not text_runs:
                raise Exception("No text content found in PowerPoint slides")
            
            # Analyze text
            
            analysis = await self.translation_service.analyze_material_text(text_runs)
            print("Analysis complete")
    
            self.update_progress(0.9)  # 90% done
    
            # Insert notes
            print("Inserting notes")
            self.supabase.table("lecture_materials").update({
                "notes": analysis
            }).eq("material_id", self.lecture_material_id).execute()
    
            # Mark done
            self.update_progress(1.0)  # 100% done
            print(f"PowerPoint processing completed for material {self.lecture_material_id}")
            
        except Exception as e:
            # If something fails, update the DB accordingly
            error_msg = str(e)
            try:
                # Make sure we're updating the correct table with sanitized error message
                self.supabase.table("lecture_materials").update({
                    "error": sanitize_text(error_msg),
                    "progress": 0.0  # Reset progress
                }).eq("material_id", self.lecture_material_id).execute()
            except Exception as db_error:
                print(f"Error updating database with error status: {db_error}")
                
            print(f"Error processing PowerPoint material {self.lecture_material_id}: {error_msg}")
            raise e
        finally:
            # Clean up temporary file
            try:
                if os.path.exists(self.file_path):
                    os.remove(self.file_path)
                    print(f"Temporary file removed: {self.file_path}")
            except Exception as cleanup_error:
                print(f"Error cleaning up temporary file: {cleanup_error}")

    async def process_pdf(self):
        """Process the PDF and update 'progress' column as each step completes."""
        try:

            # 1) Extract text from PDF
            text_content = extract_text_from_pdf()
            if not text_content:
                raise Exception("Could not extract text from PDF")
            self.update_progress(0.2)  # 20% done

            # 2) Split text into paragraphs
            paragraphs = split_into_paragraphs(text_content)
            self.supabase.table("lecture_materials").update({
                "paragraphs": paragraphs
            }).eq("material_id", self.lecture_material_id).execute()
            self.update_progress(0.3)  # 30% done
            
            # 3) Analyze text content
            
            analysis = await self.translation_service.analyze_material_text(paragraphs)
            print("Analysis")
            print(analysis)
            self.update_progress(0.5)  # 50% done

            # 4) Insert notes
            print("Inserting notes")
            self.supabase.table("lecture_materials").update({
                "notes": analysis
            }).eq("material_id", self.lecture_material_id).execute()
            self.update_progress(0.9)  # 90% done
            
            # 11) Mark done
            self.update_progress(1.0)  # 100% done
            print(f"PDF processing completed for lecture {self.lecture_material_id}")

        except Exception as e:
            # If something fails, update the DB accordingly
            error_msg = str(e)
            try:
                # Note: Changed from "lectures" to "lecture_materials"
                self.supabase.table("lecture_materials").update({
                    "error": sanitize_text(error_msg),
                    "progress": 0.0
                }).eq("material_id", self.lecture_material_id).execute()
            except Exception as db_error:
                print(f"Error updating database with error status: {db_error}")
                
            print(f"Error processing PDF material {self.lecture_material_id}: {error_msg}")
            raise e
        finally:
            # Clean up temporary file
            if os.path.exists(self.file_path):
                os.remove(self.file_path)

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file using PyPDF2."""
    # Extract text from a PDF file
    try:
        print(f"Extracting text from PDF: {file_path}")
        print(f"File size: {os.path.getsize(file_path)} bytes")
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            print(f"Number of pages: {len(pdf_reader.pages)}")
            # Extract text from each page
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += sanitize_text(page.extract_text()) + "\n\n"
            return text.strip()
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return ""

def split_into_paragraphs(text: str) -> List[str]:
    """Split text into logical paragraphs."""
    # First split by double newlines which typically indicate paragraphs
    initial_split = [p.strip() for p in text.split('\n\n') if p.strip()]
    paragraphs = []
    
    for block in initial_split:
        # Further split long blocks if they contain single newlines
        if len(block) > 500 and '\n' in block:
            sub_paragraphs = [p.strip() for p in block.split('\n') if p.strip()]
            paragraphs.extend(sub_paragraphs)
        else:
            paragraphs.append(block)
    return paragraphs

def sanitize_text(text: str) -> str:
    """Sanitize text by removing unwanted characters."""
    unwanted_chars = ['\u0000', '\n', '\r']

    for char in unwanted_chars:
        text = text.replace(char, '')
    return text